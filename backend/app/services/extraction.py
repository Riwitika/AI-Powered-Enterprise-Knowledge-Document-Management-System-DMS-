import os
import logging
from PIL import Image

logger = logging.getLogger(__name__)

# Lazy loading of format-specific libraries to ensure the backend starts even if some package has import issues
def extract_text_from_pdf(file_path: str) -> str:
    import fitz  # PyMuPDF
    text = ""
    try:
        with fitz.open(file_path) as doc:
            for page in doc:
                text += page.get_text() + "\n"
    except Exception as e:
        logger.error(f"Error parsing PDF {file_path}: {e}")
    return text

def extract_text_from_docx(file_path: str) -> str:
    import docx
    text = []
    try:
        doc = docx.Document(file_path)
        for para in doc.paragraphs:
            if para.text:
                text.append(para.text)
    except Exception as e:
        logger.error(f"Error parsing DOCX {file_path}: {e}")
    return "\n".join(text)

def extract_text_from_pptx(file_path: str) -> str:
    from pptx import Presentation
    text = []
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text.append(shape.text)
    except Exception as e:
        logger.error(f"Error parsing PPTX {file_path}: {e}")
    return "\n".join(text)

def extract_text_from_xlsx(file_path: str) -> str:
    import openpyxl
    text = []
    try:
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        for sheet in wb.worksheets:
            text.append(f"--- Sheet: {sheet.title} ---")
            for row in sheet.iter_rows(values_only=True):
                row_str = " | ".join([str(cell) for cell in row if cell is not None])
                if row_str:
                    text.append(row_str)
    except Exception as e:
        logger.error(f"Error parsing XLSX {file_path}: {e}")
    return "\n".join(text)

def extract_text_from_image(file_path: str) -> str:
    import pytesseract
    text = ""
    try:
        image = Image.open(file_path)
        text = pytesseract.image_to_string(image)
    except Exception as e:
        logger.error(f"Error performing OCR on {file_path}: {e}")
        text = f"[OCR Failed or Tesseract not installed. Filename: {os.path.basename(file_path)}]"
    return text

def extract_text_from_txt(file_path: str) -> str:
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception as e:
        logger.error(f"Error reading TXT {file_path}: {e}")
        return ""

def extract_document_text(file_path: str, file_type: str) -> str:
    file_type = file_type.lower().strip(".")
    if file_type == "pdf":
        return extract_text_from_pdf(file_path)
    elif file_type in ["docx", "doc"]:
        return extract_text_from_docx(file_path)
    elif file_type in ["pptx", "ppt"]:
        return extract_text_from_pptx(file_path)
    elif file_type in ["xlsx", "xls"]:
        return extract_text_from_xlsx(file_path)
    elif file_type in ["png", "jpg", "jpeg", "tiff", "bmp", "gif"]:
        return extract_text_from_image(file_path)
    elif file_type in ["txt", "md", "csv"]:
        return extract_text_from_txt(file_path)
    else:
        # Fallback to general text read
        return extract_text_from_txt(file_path)
